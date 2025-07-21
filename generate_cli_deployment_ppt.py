from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_title_slide(prs):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "CLI Deployment Tool"
    subtitle.text = "One-Command Cloud Deployments for Front-End Teams\n<Your Name>, Senior Product Developer"

    # Add background placeholder note
    tx_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.0), Inches(9), Inches(1))
    tf = tx_box.text_frame
    tf.text = "[Hero background image – replace with code→cloud graphic]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)


def add_content_slide(prs, title_text, bullet_lines):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for line in bullet_lines:
        p = body.add_paragraph()
        p.text = line
        p.font.size = Pt(22)
        p.level = 0


def add_meme_slide(prs, caption):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    # Placeholder rectangle for meme
    slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4)).text_frame.text = "[Insert MEME here]"
    cap_box = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(8), Inches(1))
    cap_tf = cap_box.text_frame
    cap_tf.text = caption
    cap_tf.paragraphs[0].font.size = Pt(20)
    cap_tf.paragraphs[0].font.bold = True
    cap_tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_table_slide(prs, title_text, table_data):
    rows = len(table_data)
    cols = len(table_data[0])
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text

    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
    # Table styling
    for i, row in enumerate(table_data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(14)
            para.font.bold = True if i == 0 else False
            para.alignment = PP_ALIGN.CENTER


def generate_presentation(path="CLI_Deployment_Tool_Presentation.pptx"):
    prs = Presentation()

    # 1. Title Slide
    add_title_slide(prs)

    # 2. Why & Target Audience
    add_content_slide(prs, "Why It Matters & Target Audience", [
        "Student & indie dev teams struggle to ship side-projects live.",
        "Target: Computer-science students & hackathon teams.",
        "Market stat: 78% of CS grads deploy ≤3 personal projects live."
    ])

    # 3. Vision Slide
    add_content_slide(prs, "Our Vision", [
        "From Laptop to Production in <5 min.",
        "Simplicity · Affordability · Observability"
    ])

    # 4. Meme slide
    add_meme_slide(prs, "When you finally SSH into prod and nothing works 🙃")

    # 5. Product Demo Flow
    add_content_slide(prs, "Product Demo Flow", [
        "Code → deploy-tool init → S3 artifact → Docker image → Terraform apply → EC2 → Monitoring URLs"
    ])

    # 6. Key Business Benefits
    add_content_slide(prs, "Key Business Benefits", [
        "Cuts deployment time by 90%.",
        "Zero vendor lock-in – runs on any AWS account.",
        "Built-in dashboards = faster insights → happier users."
    ])

    # 7. Revenue / Adoption Model
    add_content_slide(prs, "Revenue & Adoption Model", [
        "OSS core, paid ‘Pro Pack’ for multiservice deploys & secret vault.",
        "Upsell to universities as DevOps curriculum add-on."
    ])

    # 8. Competitive Landscape
    table_data = [
        ["Feature", "CLI Deployment Tool", "Render", "Vercel", "Netlify"],
        ["Full AWS Control", "✅", "❌", "❌", "❌"],
        ["Terraform infra code", "✅", "❌", "❌", "❌"],
        ["Rollback CLI", "✅", "⚠️", "✅", "⚠️"],
    ]
    add_table_slide(prs, "Competitive Landscape", table_data)

    # 9. Technical Architecture Deep Dive
    add_content_slide(prs, "Technical Architecture Deep Dive", [
        "GitPython → Build (Node, Vite) → Docker Hub push → Terraform modules → EC2", 
        "Docker Compose: app + Prometheus + Grafana + Node Exporter + Blackbox Exporter",
        "Auto-outputs public URL w/ security groups configured"
    ])

    #10. Code Snippet Showcase
    add_content_slide(prs, "Code Snippet Showcase", [
        "infra.py → call: subprocess.run(['terraform', 'apply', '-auto-approve'])"
    ])

    #11. Monitoring Stack Explained
    add_content_slide(prs, "Monitoring Stack Explained", [
        "Prometheus scraping app metrics (port 9090).",
        "Grafana dashboards (port 3000, admin/admin).",
        "[MEME – When Grafana hits 100% uptime and you screenshot it like 🤳]"
    ])

    #12. Risk & Mitigation
    add_content_slide(prs, "Risk & Mitigation", [
        "Drawbacks: hard-coded creds in demo, AWS costs after free tier, requires Docker daemon.",
        "Mitigation roadmap: secret manager, cost alerts, container-native build kit."
    ])

    #13. Go-to-Market Timeline
    add_content_slide(prs, "Go-to-Market Timeline", [
        "Q1: Beta at 3 universities.",
        "Q2: GitHub Marketplace launch.",
        "Q3: Pro Pack SaaS."
    ])

    #14. Funding Ask
    add_content_slide(prs, "Funding Ask", [
        "Seeking $300k pre-seed → runway 18 months.",
        "Use of funds: 50% engineering, 30% community building, 20% infra credits.",
        "Expected ARR: $1.2M by Year 3."
    ])

    #15. Closing & Call to Action
    add_content_slide(prs, "It Just Works 🚀", [
        "GitHub repo | Demo video | contact@email.dev",
        "[Insert MEME – ‘Invest now, thank us later’]"
    ])

    prs.save(path)
    print(f"Presentation saved to {path}")


if __name__ == "__main__":
    generate_presentation()